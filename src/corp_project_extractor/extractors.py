"""Per-filetype text extractors. Each writes _extracted/<slugified_rel_path>.md
with a YAML provenance header.

Supported: .pptx (python-pptx), .pdf (pdfplumber), .docx (python-docx),
           .xlsx/.xls (openpyxl), .csv (stdlib csv)
"""
from __future__ import annotations

import csv
import logging
import re
from datetime import datetime
from pathlib import Path
from typing import Optional

import yaml

from corp_project_extractor.config import get_settings
from corp_project_extractor.models import ExtractionResult

log = logging.getLogger(__name__)


# ── Path helpers ───────────────────────────────────────────────────────────────

def _slug(text: str) -> str:
    """Replace spaces and path separators with underscores, keep safe chars."""
    return re.sub(r"[^\w\-.]", "_", text)


def get_extracted_path(project_path: Path, file_path: Path) -> Path:
    """Return destination .md path under _extracted/, preserving relative structure."""
    settings = get_settings()
    extracted_dir = project_path / settings.extracted_dir
    try:
        rel = file_path.relative_to(project_path)
    except ValueError:
        rel = Path(file_path.name)
    # Mirror directory structure, append .md to filename
    out = extracted_dir / rel.parent / (file_path.name + ".md")
    return out


# ── Provenance header ──────────────────────────────────────────────────────────

def _write_extracted(
    out_path: Path,
    source_rel: str,
    source_hash: str,
    category: str,
    doc_role: str,
    extractor: str,
    metadata: dict,
    body: str,
) -> int:
    """Write markdown file with YAML front-matter. Returns word count."""
    out_path.parent.mkdir(parents=True, exist_ok=True)
    header = {
        "source": source_rel.replace("\\", "/"),
        "source_hash": source_hash,
        "category": category,
        "doc_role": doc_role,
        "extractor": extractor,
        "extracted_at": datetime.now().isoformat(timespec="seconds"),
        **metadata,
    }
    header_yaml = yaml.dump(header, allow_unicode=True, default_flow_style=False).strip()
    content = f"---\n{header_yaml}\n---\n\n{body}"
    out_path.write_text(content, encoding="utf-8")
    return len(body.split())


def _cell_str(value, max_chars: int = 200) -> str:
    if value is None:
        return ""
    s = str(value).replace("\n", " ").replace("\r", " ").replace("|", "\\|")
    return s[:max_chars] + ("…" if len(s) > max_chars else "")


# ── PPTX ──────────────────────────────────────────────────────────────────────

def extract_pptx(
    file_path: Path,
    project_path: Path,
    category: str,
    doc_role: str,
    source_hash: str = "",
) -> ExtractionResult:
    from pptx import Presentation  # type: ignore

    settings = get_settings()
    out_path = get_extracted_path(project_path, file_path)

    try:
        prs = Presentation(str(file_path))
    except Exception as e:
        return ExtractionResult(success=False, output_path=None, word_count=0, error=str(e))

    lines: list[str] = [f"# {file_path.stem}", ""]
    slide_count = len(prs.slides)
    extracted_slides = 0

    for i, slide in enumerate(prs.slides, 1):
        title_text = ""
        content_lines: list[str] = []
        notes_text = ""

        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            ph = getattr(shape, "placeholder_format", None)
            ph_idx = ph.idx if ph else -1

            if ph_idx == 0:  # title placeholder
                title_text = shape.text_frame.text.strip()
                continue

            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if not text:
                    continue
                indent = "  " * para.level
                content_lines.append(f"{indent}- {text}")

        if settings.extraction.pptx_include_notes and slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            if notes_tf:
                notes_text = notes_tf.text.strip()

        if not title_text and not content_lines and not notes_text:
            continue  # skip empty slides

        extracted_slides += 1
        slide_label = f"Slide {i}" if settings.extraction.pptx_include_slide_numbers else ""
        heading = f"## {slide_label}: {title_text}" if title_text else f"## {slide_label}"
        lines.append(heading.strip(": "))
        if content_lines:
            lines.extend(content_lines)
        if notes_text:
            lines.append("")
            lines.append(f"> **Notes:** {notes_text}")
        lines.append("")

    body = "\n".join(lines)
    source_rel = str(file_path.relative_to(project_path))
    wc = _write_extracted(
        out_path, source_rel, source_hash, category, doc_role, "pptx",
        {"slides": slide_count, "extracted_slides": extracted_slides}, body,
    )
    return ExtractionResult(success=True, output_path=out_path, word_count=wc)


# ── PDF ───────────────────────────────────────────────────────────────────────

def extract_pdf(
    file_path: Path,
    project_path: Path,
    category: str,
    doc_role: str,
    source_hash: str = "",
) -> ExtractionResult:
    import pdfplumber  # type: ignore

    out_path = get_extracted_path(project_path, file_path)

    try:
        pdf = pdfplumber.open(str(file_path))
    except Exception as e:
        return ExtractionResult(success=False, output_path=None, word_count=0, error=str(e))

    lines: list[str] = [f"# {file_path.stem}", ""]
    page_count = len(pdf.pages)
    extracted_pages = 0

    with pdf:
        for i, page in enumerate(pdf.pages, 1):
            try:
                text = page.extract_text()
            except Exception as e:
                log.warning("PDF page %d extraction error in %s: %s", i, file_path.name, e)
                continue
            if text and text.strip():
                lines.append(f"## Page {i}")
                lines.append(text.strip())
                lines.append("")
                extracted_pages += 1

    body = "\n".join(lines)
    source_rel = str(file_path.relative_to(project_path))
    wc = _write_extracted(
        out_path, source_rel, source_hash, category, doc_role, "pdf",
        {"pages": page_count, "extracted_pages": extracted_pages}, body,
    )
    return ExtractionResult(success=True, output_path=out_path, word_count=wc)


# ── DOCX ──────────────────────────────────────────────────────────────────────

def extract_docx(
    file_path: Path,
    project_path: Path,
    category: str,
    doc_role: str,
    source_hash: str = "",
) -> ExtractionResult:
    from docx import Document  # type: ignore

    out_path = get_extracted_path(project_path, file_path)

    try:
        doc = Document(str(file_path))
    except Exception as e:
        return ExtractionResult(success=False, output_path=None, word_count=0, error=str(e))

    lines: list[str] = [f"# {file_path.stem}", ""]
    para_count = 0

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue
        style = para.style.name if para.style else ""
        if style.startswith("Heading 1"):
            lines.append(f"## {text}")
        elif style.startswith("Heading 2"):
            lines.append(f"### {text}")
        elif style.startswith("Heading 3"):
            lines.append(f"#### {text}")
        elif style.startswith("Heading"):
            lines.append(f"##### {text}")
        elif style.startswith("List"):
            lines.append(f"- {text}")
        else:
            lines.append(text)
        para_count += 1

    # Also extract tables
    for table in doc.tables:
        lines.append("")
        # Header row
        header = [_cell_str(c.text) for c in table.rows[0].cells] if table.rows else []
        if header:
            lines.append("| " + " | ".join(header) + " |")
            lines.append("| " + " | ".join(["---"] * len(header)) + " |")
            for row in table.rows[1:]:
                cells = [_cell_str(c.text) for c in row.cells]
                lines.append("| " + " | ".join(cells) + " |")
        lines.append("")

    body = "\n".join(lines)
    source_rel = str(file_path.relative_to(project_path))
    wc = _write_extracted(
        out_path, source_rel, source_hash, category, doc_role, "docx",
        {"paragraphs": para_count}, body,
    )
    return ExtractionResult(success=True, output_path=out_path, word_count=wc)


# ── XLSX ──────────────────────────────────────────────────────────────────────

def _is_questionnaire_sheet(sheet_name: str) -> bool:
    """True if this sheet should use full-row questionnaire extraction."""
    settings = get_settings()
    name_lower = sheet_name.lower()
    return any(m.lower() in name_lower for m in settings.extraction.xlsx_questionnaire_markers)


def extract_xlsx(
    file_path: Path,
    project_path: Path,
    category: str,
    doc_role: str,
    source_hash: str = "",
) -> ExtractionResult:
    from openpyxl import load_workbook  # type: ignore

    settings = get_settings()
    out_path = get_extracted_path(project_path, file_path)
    trunc = settings.extraction.xlsx_cell_truncate_chars

    try:
        wb = load_workbook(str(file_path), read_only=True, data_only=True)
    except Exception as e:
        return ExtractionResult(success=False, output_path=None, word_count=0, error=str(e))

    lines: list[str] = [f"# {file_path.stem}", ""]
    sheet_count = len(wb.sheetnames)

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        lines.append(f"## Sheet: {sheet_name}")
        lines.append("")

        # Materialise all rows (read_only iterator is one-pass)
        try:
            all_rows = list(ws.iter_rows(values_only=True))
        except Exception as e:
            lines.append(f"_(error reading sheet: {e})_")
            lines.append("")
            continue

        if not all_rows:
            lines.append("_(empty sheet)_")
            lines.append("")
            continue

        # Find first non-empty row as header
        headers: list[str] = []
        data_start = 0
        for idx, row in enumerate(all_rows):
            if any(c is not None for c in row):
                headers = [_cell_str(c, trunc) for c in row]
                data_start = idx + 1
                break

        if not headers:
            lines.append("_(no data)_")
            lines.append("")
            continue

        questionnaire_mode = _is_questionnaire_sheet(sheet_name)
        data_rows = all_rows[data_start:]
        non_empty = [r for r in data_rows if any(c is not None for c in r)]

        if questionnaire_mode:
            max_rows = settings.extraction.xlsx_questionnaire_max_rows
            lines.append(f"_Mode: questionnaire — {len(non_empty)} rows_")
            lines.append("")
            lines.append("| " + " | ".join(headers) + " |")
            lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
            for row in non_empty[:max_rows]:
                cells = [_cell_str(c, trunc) for c in row]
                # Pad/trim to match header width
                cells = cells[:len(headers)] + [""] * (len(headers) - len(cells))
                lines.append("| " + " | ".join(cells) + " |")
            if len(non_empty) > max_rows:
                lines.append(f"_… {len(non_empty) - max_rows} more rows_")
        else:
            max_rows = settings.extraction.xlsx_data_mode_max_rows
            total = len(non_empty)
            lines.append(f"_Mode: data — {total} rows, showing first {min(total, max_rows)}_")
            lines.append("")
            lines.append("**Columns:** " + ", ".join(f"`{h}`" for h in headers if h))
            lines.append("")
            sample = non_empty[:max_rows]
            if sample:
                lines.append("| " + " | ".join(headers) + " |")
                lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
                for row in sample:
                    cells = [_cell_str(c, trunc) for c in row]
                    cells = cells[:len(headers)] + [""] * (len(headers) - len(cells))
                    lines.append("| " + " | ".join(cells) + " |")
                if total > max_rows:
                    lines.append(f"_… {total - max_rows} more rows_")
        lines.append("")

    wb.close()
    body = "\n".join(lines)
    source_rel = str(file_path.relative_to(project_path))
    wc = _write_extracted(
        out_path, source_rel, source_hash, category, doc_role, "xlsx",
        {"sheets": sheet_count}, body,
    )
    return ExtractionResult(success=True, output_path=out_path, word_count=wc)


# ── CSV ───────────────────────────────────────────────────────────────────────

def extract_csv(
    file_path: Path,
    project_path: Path,
    category: str,
    doc_role: str,
    source_hash: str = "",
) -> ExtractionResult:
    out_path = get_extracted_path(project_path, file_path)
    rows: list[list[str]] = []

    for encoding in ("utf-8-sig", "utf-8", "latin-1"):
        try:
            with open(file_path, encoding=encoding, newline="") as f:
                rows = list(csv.reader(f))
            break
        except (UnicodeDecodeError, Exception):
            continue

    headers = rows[0] if rows else []
    data_rows = rows[1:]
    total = len(data_rows)

    lines: list[str] = [f"# {file_path.stem}", ""]
    lines.append(f"**Rows:** {total}  |  **Columns:** {len(headers)}")
    lines.append("")
    if headers:
        lines.append("**Headers:** " + ", ".join(f"`{h}`" for h in headers))
        lines.append("")
    sample = data_rows[:10]
    if sample and headers:
        lines.append("**Sample rows (first 10):**")
        lines.append("")
        lines.append("| " + " | ".join(headers) + " |")
        lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
        for row in sample:
            cells = [c.replace("|", "\\|") for c in row]
            lines.append("| " + " | ".join(cells) + " |")

    body = "\n".join(lines)
    source_rel = str(file_path.relative_to(project_path))
    wc = _write_extracted(
        out_path, source_rel, source_hash, category, doc_role, "csv",
        {"rows": total, "columns": len(headers)}, body,
    )
    return ExtractionResult(success=True, output_path=out_path, word_count=wc)


# ── Dispatcher ────────────────────────────────────────────────────────────────

def extract(
    file_path: Path,
    project_path: Path,
    category: str,
    doc_role: str,
    source_hash: str = "",
) -> ExtractionResult:
    """Dispatch extraction to the right extractor. Returns ExtractionResult."""
    settings = get_settings()
    ext = file_path.suffix.lower()

    # Size check
    try:
        size_mb = file_path.stat().st_size / (1024 * 1024)
    except OSError as e:
        return ExtractionResult(False, None, 0, f"Cannot stat file (cloud-only?): {e}")
    if size_mb > settings.extraction.max_file_size_mb:
        return ExtractionResult(False, None, 0,
                                f"File too large ({size_mb:.0f} MB > {settings.extraction.max_file_size_mb} MB)")

    # Skip list
    if ext in settings.skip_extensions:
        return ExtractionResult(False, None, 0, f"skipped extension: {ext}")

    try:
        if ext == ".pptx":
            return extract_pptx(file_path, project_path, category, doc_role, source_hash)
        elif ext == ".pdf":
            return extract_pdf(file_path, project_path, category, doc_role, source_hash)
        elif ext == ".docx":
            return extract_docx(file_path, project_path, category, doc_role, source_hash)
        elif ext in (".xlsx", ".xls"):
            return extract_xlsx(file_path, project_path, category, doc_role, source_hash)
        elif ext == ".csv":
            return extract_csv(file_path, project_path, category, doc_role, source_hash)
        else:
            return ExtractionResult(False, None, 0, f"unsupported extension: {ext}")
    except Exception as e:
        log.warning("Extraction failed for %s: %s", file_path.name, e, exc_info=True)
        return ExtractionResult(False, None, 0, str(e))
